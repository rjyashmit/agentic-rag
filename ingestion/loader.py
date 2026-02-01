import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd

def load_document(path):
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        reader = PdfReader(path)
        return "\n".join(p.extract_text() for p in reader.pages)

    elif ext == ".docx":
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)

    elif ext == ".pptx":
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    elif ext == ".xlsx":
        df = pd.read_excel(path)
        return df.to_string()

    elif ext == ".txt":
        with open(path, "r", encoding="utf-8") as f:
            return f.read()

    else:
        raise ValueError("Unsupported file type")